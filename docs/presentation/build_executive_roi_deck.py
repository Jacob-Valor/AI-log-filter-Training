from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT_PATH = Path("docs/presentation/AI-Log-Filter-Executive-ROI.pptx")
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# Fictional-grid-inspired palette (adapted for executive cybersecurity content)
BG_LIGHT = RGBColor(249, 245, 239)
BG_LIGHT_ALT = RGBColor(244, 238, 231)
BG_DARK = RGBColor(37, 33, 52)
GRID_LIGHT = RGBColor(227, 215, 203)
GRID_DARK = RGBColor(88, 82, 110)
TEXT_DARK = RGBColor(34, 36, 48)
TEXT_MID = RGBColor(89, 83, 108)
TEXT_LIGHT = RGBColor(240, 236, 246)
ACCENT_RED = RGBColor(232, 74, 98)
ACCENT_CORAL = RGBColor(250, 120, 96)
ACCENT_PINK = RGBColor(239, 159, 177)
ACCENT_BLUE = RGBColor(78, 118, 178)
ACCENT_GREEN = RGBColor(47, 146, 103)
WHITE = RGBColor(255, 255, 255)

TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"


def add_full_rect(slide, color: RGBColor):
    rect = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        SLIDE_W,
        SLIDE_H,
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect


def add_grid(slide, dark: bool):
    slide_w = SLIDE_W
    slide_h = SLIDE_H
    line_color = GRID_DARK if dark else GRID_LIGHT

    x = 0.0
    while x <= 13.333:
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(0), Inches(x), slide_h
        )
        line.line.color.rgb = line_color
        line.line.width = Pt(0.8)
        line.line.transparency = 0.8
        x += 0.72

    y = 0.0
    while y <= 7.5:
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(y), slide_w, Inches(y)
        )
        line.line.color.rgb = line_color
        line.line.width = Pt(0.8)
        line.line.transparency = 0.8
        y += 0.72


def add_decor(slide, dark: bool):
    if dark:
        corner = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, Inches(10.6), Inches(-1.2), Inches(4.8), Inches(4.8)
        )
        corner.fill.solid()
        corner.fill.fore_color.rgb = ACCENT_PINK
        corner.fill.transparency = 0.72
        corner.line.fill.background()

        strip = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.CHEVRON,
            Inches(-0.2),
            Inches(5.9),
            Inches(4.2),
            Inches(1.7),
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = ACCENT_RED
        strip.fill.transparency = 0.18
        strip.line.fill.background()
    else:
        blob_a = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(-0.9),
            Inches(5.5),
            Inches(3.0),
            Inches(3.0),
        )
        blob_a.fill.solid()
        blob_a.fill.fore_color.rgb = ACCENT_PINK
        blob_a.fill.transparency = 0.55
        blob_a.line.fill.background()

        blob_b = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(11.0),
            Inches(-0.7),
            Inches(3.2),
            Inches(3.2),
        )
        blob_b.fill.solid()
        blob_b.fill.fore_color.rgb = ACCENT_CORAL
        blob_b.fill.transparency = 0.72
        blob_b.line.fill.background()

        stripe = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE,
            Inches(9.7),
            Inches(6.5),
            Inches(3.9),
            Inches(1.0),
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = ACCENT_RED
        stripe.fill.transparency = 0.75
        stripe.line.fill.background()


def apply_background(slide, mode: str):
    if mode == "dark":
        add_full_rect(slide, BG_DARK)
        add_grid(slide, dark=True)
        add_decor(slide, dark=True)
    elif mode == "light_alt":
        add_full_rect(slide, BG_LIGHT_ALT)
        add_grid(slide, dark=False)
        add_decor(slide, dark=False)
    else:
        add_full_rect(slide, BG_LIGHT)
        add_grid(slide, dark=False)
        add_decor(slide, dark=False)


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=20,
    bold=False,
    color=TEXT_DARK,
    align=PP_ALIGN.LEFT,
    font=BODY_FONT,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1.06
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size=17, color=TEXT_DARK):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(3)
    tf.margin_bottom = Pt(3)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        p.line_spacing = 1.08
        p.font.name = BODY_FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def add_tag(slide, text, dark=False):
    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(0.35), Inches(2.55), Inches(0.5)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = ACCENT_RED
    tag.line.fill.background()

    tf = tag.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.alignment = PP_ALIGN.CENTER
    p.font.name = BODY_FONT
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_card(slide, left, top, width, height, title, value, accent):
    shadow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left + 0.06),
        Inches(top + 0.06),
        Inches(width),
        Inches(height),
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(192, 179, 169)
    shadow.fill.transparency = 0.7
    shadow.line.fill.background()

    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = accent
    card.line.width = Pt(2)

    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.text = title
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = BODY_FONT
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_MID

    p2 = tf.add_paragraph()
    p2.text = value
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = TITLE_FONT
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = accent


def slide_title(slide, text, dark=False):
    add_text_box(
        slide,
        0.65,
        0.92,
        11.8,
        0.85,
        text,
        size=33,
        bold=True,
        color=TEXT_LIGHT if dark else TEXT_DARK,
        font=TITLE_FONT,
    )


def build_deck(path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # 1. Cover
    s = prs.slides.add_slide(blank)
    apply_background(s, "dark")
    add_tag(s, "Executive Brief", dark=True)
    slide_title(s, "AI-Driven Log Filtering for SIEM Efficiency", dark=True)
    add_text_box(
        s,
        0.68,
        1.95,
        10.8,
        0.7,
        "Executive ROI Review | IBM QRadar Focus | February 2026",
        size=19,
        color=TEXT_LIGHT,
        font=BODY_FONT,
    )
    add_card(s, 0.75, 4.3, 3.9, 1.55, "SIEM Volume", "40-60%", ACCENT_RED)
    add_card(s, 4.95, 4.3, 3.9, 1.55, "Critical Recall Target", ">99.5%", ACCENT_BLUE)
    add_card(s, 9.15, 4.3, 3.35, 1.55, "Safety Mode", "Fail-open", ACCENT_GREEN)
    add_text_box(
        s,
        0.7,
        6.78,
        12.2,
        0.45,
        "Source basis: repository benchmark, architecture, tracking, and validation artifacts.",
        size=11,
        color=RGBColor(205, 197, 219),
    )

    # 2. Snapshot
    s = prs.slides.add_slide(blank)
    apply_background(s, "light")
    add_tag(s, "At a glance")
    slide_title(s, "Executive Snapshot")

    panel = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.75), Inches(7.35), Inches(4.9)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = RGBColor(216, 198, 184)
    panel.line.width = Pt(1.5)
    add_bullets(
        s,
        1.02,
        2.05,
        6.85,
        4.4,
        [
            "AI pre-filtering controls SIEM ingest volume before QRadar licensing thresholds are hit.",
            "Zero-loss safety posture: fail-open fallback + immutable decision audit trail.",
            "January benchmark package exceeds latency, throughput, and critical recall targets.",
            "Financial model indicates six-figure annual savings at enterprise ingest volumes.",
        ],
        size=18,
    )
    add_card(s, 8.35, 1.75, 2.15, 1.95, "EPS Reduction", "79%", ACCENT_RED)
    add_card(s, 10.68, 1.75, 2.15, 1.95, "Recall", "99.7%", ACCENT_BLUE)
    add_card(s, 8.35, 3.95, 2.15, 1.95, "P95 Latency", "42 ms", ACCENT_CORAL)
    add_card(s, 10.68, 3.95, 2.15, 1.95, "Annual Savings", "$180k-$360k", ACCENT_GREEN)

    # 3. Problem + opportunity
    s = prs.slides.add_slide(blank)
    apply_background(s, "light_alt")
    add_tag(s, "Business case")
    slide_title(s, "Problem and Opportunity")
    add_bullets(
        s,
        0.82,
        1.92,
        5.9,
        4.8,
        [
            "Raw SIEM ingestion is expensive and creates analyst overload.",
            "Most logs are routine/noise while critical detections remain sparse.",
            "The system compresses volume while preserving high-severity signal quality.",
            "Benchmark scenario: 15,000 input EPS reduced to 3,150 QRadar EPS.",
        ],
        size=18,
    )
    chart_data = CategoryChartData()
    chart_data.categories = ["Before filter", "After filter"]
    chart_data.add_series("QRadar EPS", (15000, 3150))
    chart = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(6.7),
        Inches(1.95),
        Inches(5.9),
        Inches(4.15),
        chart_data,
    ).chart
    chart.has_legend = False
    chart.value_axis.maximum_scale = 16000
    chart.value_axis.minimum_scale = 0
    chart.value_axis.major_unit = 4000
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.tick_labels.font.size = Pt(12)
    chart.value_axis.tick_labels.font.size = Pt(10)
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = "#,##0"
    plot.data_labels.font.size = Pt(12)
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = ACCENT_RED

    add_text_box(
        s,
        6.86,
        6.15,
        5.65,
        0.5,
        "Routing volume reduction in benchmark package: 79%",
        size=14,
        bold=True,
        color=ACCENT_RED,
        align=PP_ALIGN.CENTER,
    )

    # 4. Routing model
    s = prs.slides.add_slide(blank)
    apply_background(s, "light")
    add_tag(s, "Operating model")
    slide_title(s, "How the Platform Routes Logs")
    add_card(s, 0.72, 1.9, 2.95, 2.2, "Critical", "Immediate QRadar", ACCENT_RED)
    add_card(s, 3.97, 1.9, 2.95, 2.2, "Suspicious", "Priority QRadar", ACCENT_CORAL)
    add_card(s, 7.22, 1.9, 2.95, 2.2, "Routine", "Cold storage", ACCENT_BLUE)
    add_card(s, 10.47, 1.9, 2.15, 2.2, "Noise", "Summary + archive", ACCENT_GREEN)

    rail = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.9),
        Inches(4.65),
        Inches(11.95),
        Inches(1.55),
    )
    rail.fill.solid()
    rail.fill.fore_color.rgb = WHITE
    rail.line.color.rgb = RGBColor(213, 196, 184)
    rail.line.width = Pt(1.4)
    tf = rail.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = "Parse (5ms)  ->  Compliance Check (1ms)  ->  Enrich (10ms)  ->  Classify (20ms)  ->  Route (5ms)"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = BODY_FONT
    p1.font.size = Pt(15)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_DARK
    p2 = tf.add_paragraph()
    p2.text = "Approximate end-to-end latency: ~50ms"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = BODY_FONT
    p2.font.size = Pt(13)
    p2.font.color.rgb = TEXT_MID

    # 5. Safety/compliance
    s = prs.slides.add_slide(blank)
    apply_background(s, "dark")
    add_tag(s, "Risk controls", dark=True)
    slide_title(s, "Safety and Compliance Controls", dark=True)

    left = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(1.75),
        Inches(6.0),
        Inches(4.95),
    )
    left.fill.solid()
    left.fill.fore_color.rgb = RGBColor(52, 45, 72)
    left.line.color.rgb = ACCENT_BLUE
    left.line.width = Pt(1.4)
    add_text_box(
        s, 1.02, 2.02, 5.45, 0.5, "Core safety guarantees", size=19, bold=True, color=TEXT_LIGHT
    )
    add_bullets(
        s,
        1.02,
        2.45,
        5.35,
        4.05,
        [
            "Fail-open routing on model/system error.",
            "Circuit breaker with automated recovery cycle.",
            "Timeout caps and graceful degradation mode.",
            "Every decision logged for auditability.",
            "Cold storage retention preserves forensic evidence.",
        ],
        size=16,
        color=TEXT_LIGHT,
    )

    right = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(7.05),
        Inches(1.75),
        Inches(5.55),
        Inches(4.95),
    )
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor(56, 50, 78)
    right.line.color.rgb = ACCENT_GREEN
    right.line.width = Pt(1.4)
    add_text_box(
        s, 7.3, 2.02, 5.0, 0.5, "Compliance-first bypass", size=19, bold=True, color=TEXT_LIGHT
    )
    add_bullets(
        s,
        7.3,
        2.45,
        5.0,
        4.05,
        [
            "PCI-DSS source patterns bypass AI.",
            "HIPAA and EHR records bypass AI.",
            "SOX/financial records bypass AI.",
            "GDPR-related categories bypass AI.",
            "Regulated logs route directly to QRadar.",
        ],
        size=16,
        color=TEXT_LIGHT,
    )

    # 6. Architecture + ensemble
    s = prs.slides.add_slide(blank)
    apply_background(s, "light_alt")
    add_tag(s, "Technical backbone")
    slide_title(s, "Architecture and Model Strategy")
    add_bullets(
        s,
        0.84,
        1.9,
        6.6,
        4.65,
        [
            "Kafka ingestion for raw, classified, and audit streams.",
            "Compliance gate before ML path for regulated events.",
            "Safe ensemble combines rules, TF-IDF/XGBoost, and anomaly detector.",
            "Routing layer forwards by severity and storage policy.",
            "Prometheus + Grafana provide operational and quality visibility.",
        ],
        size=17,
    )
    pie_data = CategoryChartData()
    pie_data.categories = ["Rule-based", "TF-IDF + XGBoost", "Anomaly detector"]
    pie_data.add_series("Weight", (30, 45, 25))
    pie = s.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(7.2),
        Inches(1.85),
        Inches(5.2),
        Inches(4.3),
        pie_data,
    ).chart
    pie.has_legend = True
    pie.legend.position = XL_LEGEND_POSITION.BOTTOM
    pie.legend.include_in_layout = False
    pie.plots[0].has_data_labels = True
    pie.plots[0].data_labels.show_percentage = True
    pie.plots[0].data_labels.number_format = "0%"
    pie.series[0].points[0].format.fill.solid()
    pie.series[0].points[0].format.fill.fore_color.rgb = ACCENT_RED
    pie.series[0].points[1].format.fill.solid()
    pie.series[0].points[1].format.fill.fore_color.rgb = ACCENT_BLUE
    pie.series[0].points[2].format.fill.solid()
    pie.series[0].points[2].format.fill.fore_color.rgb = ACCENT_GREEN
    add_text_box(
        s,
        7.2,
        6.18,
        5.25,
        0.55,
        "Optional ONNX path documented with major inference and size gains.",
        size=12,
        color=TEXT_MID,
        align=PP_ALIGN.CENTER,
    )

    # 7. Performance proof
    s = prs.slides.add_slide(blank)
    apply_background(s, "light")
    add_tag(s, "Performance")
    slide_title(s, "Performance Evidence Against Targets")

    rows = [
        ("Peak EPS", "10,000", "12,847", "Exceeded"),
        ("Sustained EPS", "8,000", "9,234", "Exceeded"),
        ("P95 classification latency", "<100ms", "42ms", "Pass"),
        ("End-to-end latency", "<200ms", "89ms", "Pass"),
        ("Kafka round-trip latency", "<50ms", "23ms", "Pass"),
    ]
    tbl = s.shapes.add_table(6, 4, Inches(0.84), Inches(1.95), Inches(11.95), Inches(3.9)).table
    for idx, w in enumerate([4.25, 2.4, 2.4, 1.9]):
        tbl.columns[idx].width = Inches(w)

    headers = ["Metric", "Target", "Achieved", "Status"]
    for c, text in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_RED
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.name = BODY_FONT
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(250, 247, 242)
            for p in cell.text_frame.paragraphs:
                p.font.name = BODY_FONT
                p.font.size = Pt(12)
                p.font.color.rgb = TEXT_DARK
                p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER

    add_text_box(
        s,
        0.84,
        6.2,
        11.95,
        0.55,
        "January benchmark package demonstrates measurable headroom against operating targets.",
        size=14,
        bold=True,
        color=ACCENT_GREEN,
        align=PP_ALIGN.CENTER,
    )

    # 8. Quality proof
    s = prs.slides.add_slide(blank)
    apply_background(s, "light_alt")
    add_tag(s, "Model quality")
    slide_title(s, "Detection Quality Evidence")

    qdata = CategoryChartData()
    qdata.categories = ["Accuracy", "Critical Recall", "Critical Precision", "F1 Score"]
    qdata.add_series("Percent", (94.2, 99.7, 93.4, 93.6))
    qchart = s.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.9),
        Inches(1.95),
        Inches(6.5),
        Inches(4.3),
        qdata,
    ).chart
    qchart.has_legend = False
    qchart.value_axis.minimum_scale = 0
    qchart.value_axis.maximum_scale = 100
    qchart.value_axis.major_unit = 20
    qchart.plots[0].has_data_labels = True
    qchart.plots[0].data_labels.number_format = "0.0"
    qchart.plots[0].series[0].format.fill.solid()
    qchart.plots[0].series[0].format.fill.fore_color.rgb = ACCENT_BLUE

    right_panel = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(7.55),
        Inches(1.95),
        Inches(5.05),
        Inches(4.3),
    )
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = WHITE
    right_panel.line.color.rgb = RGBColor(214, 198, 183)
    add_bullets(
        s,
        7.82,
        2.2,
        4.55,
        3.8,
        [
            "Critical recall benchmark result: 99.7% (target >99.5%).",
            "Critical precision benchmark result: 93.4%.",
            "False negatives tightly controlled in January evidence package.",
            "Model artifact checks: 11/11 passed in validation summary.",
            "129-test suite supports release confidence.",
        ],
        size=15,
    )

    # 9. ROI slide
    s = prs.slides.add_slide(blank)
    apply_background(s, "dark")
    add_tag(s, "Financials", dark=True)
    slide_title(s, "Financial Impact and ROI", dark=True)
    add_card(s, 0.85, 1.95, 3.85, 2.0, "Baseline Tier", "$450k/yr @ 15k EPS", ACCENT_RED)
    add_card(s, 4.95, 1.95, 3.85, 2.0, "Post-filter Scenario", "~50% EPS reduction", ACCENT_BLUE)
    add_card(s, 9.05, 1.95, 3.15, 2.0, "Savings Range", "$180k-$360k/yr", ACCENT_GREEN)
    add_bullets(
        s,
        0.95,
        4.35,
        11.5,
        2.2,
        [
            "Cost-tracking reference scenario estimates ~$226k annual savings.",
            "Illustrative break-even point in the same scenario: ~96 days.",
            "Modeled infrastructure cost in reference architecture: ~$60k/year.",
            "ROI grows with sustained EPS reduction and controlled rollout execution.",
        ],
        size=17,
        color=TEXT_LIGHT,
    )

    # 10. Risk and mitigation
    s = prs.slides.add_slide(blank)
    apply_background(s, "light")
    add_tag(s, "Governance")
    slide_title(s, "Current Risk Snapshot and Mitigation")

    risk = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.85),
        Inches(1.9),
        Inches(5.8),
        Inches(4.9),
    )
    risk.fill.solid()
    risk.fill.fore_color.rgb = RGBColor(255, 242, 240)
    risk.line.color.rgb = RGBColor(236, 157, 150)
    add_text_box(
        s,
        1.08,
        2.15,
        5.35,
        0.45,
        "Observed in February reruns",
        size=18,
        bold=True,
        color=ACCENT_RED,
    )
    add_bullets(
        s,
        1.08,
        2.55,
        5.25,
        3.95,
        [
            "Shadow validation showed critical recall at 60.0%.",
            "Chaos stress test showed critical recall at 70.0%.",
            "Pattern indicates potential artifact or configuration drift.",
            "Recommendation: gate broad rollout until quality thresholds recover.",
        ],
        size=15,
    )

    mit = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(7.0),
        Inches(1.9),
        Inches(5.45),
        Inches(4.9),
    )
    mit.fill.solid()
    mit.fill.fore_color.rgb = RGBColor(236, 248, 240)
    mit.line.color.rgb = RGBColor(153, 205, 170)
    add_text_box(
        s, 7.25, 2.15, 5.0, 0.45, "Mitigation gates", size=18, bold=True, color=ACCENT_GREEN
    )
    add_bullets(
        s,
        7.25,
        2.55,
        4.95,
        3.95,
        [
            "Pin known-good model + environment versions.",
            "Require recall >=99.5% and FN <=5 per 1000.",
            "Require stress recall >=99.5% before promotion.",
            "Automate benchmark evidence in CI release gate.",
            "Provide weekly risk dashboard to leadership.",
        ],
        size=15,
    )

    # 11. 30-day plan
    s = prs.slides.add_slide(blank)
    apply_background(s, "light_alt")
    add_tag(s, "Execution")
    slide_title(s, "30-Day Execution Plan")

    line = s.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(1.25), Inches(3.15), Inches(11.95), Inches(3.15)
    )
    line.line.color.rgb = ACCENT_RED
    line.line.width = Pt(3)

    milestones = [
        (1.25, "Week 1", "Root-cause + baseline"),
        (4.0, "Week 2", "Retrain + calibrate"),
        (6.75, "Week 3", "Shadow + chaos gates"),
        (9.5, "Week 4", "5% pilot rollout"),
    ]
    colors = [ACCENT_RED, ACCENT_CORAL, ACCENT_BLUE, ACCENT_GREEN]
    for (x, title, desc), color in zip(milestones, colors, strict=True):
        dot = s.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(2.85), Inches(0.58), Inches(0.58)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        add_text_box(s, x - 0.25, 3.55, 2.0, 0.4, title, size=14, bold=True, color=TEXT_DARK)
        add_text_box(s, x - 0.25, 3.95, 2.4, 0.95, desc, size=13, color=TEXT_MID)

    ask = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.95),
        Inches(5.2),
        Inches(11.95),
        Inches(1.25),
    )
    ask.fill.solid()
    ask.fill.fore_color.rgb = WHITE
    ask.line.color.rgb = RGBColor(212, 194, 181)
    tf = ask.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Decision request: approve controlled pilot under strict quality gates and weekly executive reporting."
    p.alignment = PP_ALIGN.CENTER
    p.font.name = BODY_FONT
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK

    # 12. Sources
    s = prs.slides.add_slide(blank)
    apply_background(s, "dark")
    add_tag(s, "Appendix", dark=True)
    slide_title(s, "Evidence Sources", dark=True)
    add_bullets(
        s,
        0.95,
        1.95,
        11.95,
        4.85,
        [
            "README.md (problem framing, design principles, target metrics)",
            "reports/BENCHMARK_SUMMARY.md (January throughput/latency/quality evidence)",
            "docs/tracking/COST_TRACKING.md (ROI model and break-even scenario)",
            "docs/architecture/PRODUCTION_ARCHITECTURE.md (production controls and safeguards)",
            "docs/assessment/ASSESSMENT_SCORECARD.md (test/readiness scorecard)",
            "reports/shadow_validation/validation_20260218_075134.json (Feb regression signal)",
            "reports/chaos_test/chaos_test_20260218_073006.json (stress recall regression signal)",
        ],
        size=15,
        color=TEXT_LIGHT,
    )
    add_text_box(
        s,
        0.95,
        6.82,
        11.95,
        0.4,
        "Theme refreshed to a modern fictional-grid style (red/cream editorial layout).",
        size=11,
        color=RGBColor(192, 184, 208),
    )

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))


if __name__ == "__main__":
    build_deck(OUT_PATH)
    print(f"Created {OUT_PATH}")
